import sys
import os

# Support local library path (if needed)
sys.path.append(os.path.join(os.path.dirname(__file__), "libs"))

import threading
from concurrent.futures import ThreadPoolExecutor, as_completed
import fitz
from docx import Document
from pptx import Presentation
import tkinter as tk
from tkinter import filedialog, messagebox
from tkinter import ttk

# Extract text by file type
def extract_text_with_line_numbers(filepath):
    lines = []
    try:
        if filepath.lower().endswith(".pdf"):
            doc = fitz.open(filepath)
            for page in doc:
                text = page.get_text()
                lines += text.splitlines()
            doc.close()
        elif filepath.lower().endswith(".docx"):
            doc = Document(filepath)
            for para in doc.paragraphs:
                lines.append(para.text)
        elif filepath.lower().endswith(".pptx"):
            ppt = Presentation(filepath)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        lines += shape.text.splitlines()
        elif filepath.lower().endswith(".txt"):
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                lines = f.readlines()
    except Exception as e:
        print(f"❌ Failed to read {filepath}: {e}")
    return lines

# Process a single file and return matching lines
def process_file(full_path, keyword_lower):
    global search_cancelled
    if search_cancelled:
        return None
    try:
        lines = extract_text_with_line_numbers(full_path)
        matches = []
        for i, line in enumerate(lines):
            if keyword_lower in line.lower():
                matches.append((i + 1, line.strip()))

        if matches:
            return (full_path, matches)
    except Exception as e:
        print(f"❌ Error processing file {full_path}: {e}")
    return None

# Multi-threaded keyword search across files
def search_keyword_in_folder(folder, keyword, on_result=None):
    keyword_lower = keyword.lower()
    file_list = []

    for root_dir, _, files in os.walk(folder):
        for file in files:
            if file.lower().endswith((".pdf", ".docx", ".pptx", ".txt")):
                file_list.append(os.path.join(root_dir, file))

    cpu_count = os.cpu_count() or 4
    max_workers = max(2, min(cpu_count * 2, 16))
    print(f"🧠 Using {max_workers} threads for searching...")

    results = []
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = [executor.submit(process_file, path, keyword_lower) for path in file_list]
        for future in as_completed(futures):
            if search_cancelled:
                break
            result = future.result()
            if result:
                results.append(result)
                if on_result:
                    root.after(0, lambda res=result: on_result(res))
    return results

# Display one matched file result
def handle_result(result):
    file_path, matches = result
    keyword = keyword_entry.get().strip()
    if not keyword:
        return

    safe_path = os.path.abspath(file_path).replace("\\", "/")
    result_text.insert(tk.END, f" {safe_path}\n", "greenbar")

    for _, line_content in matches:
        result_text.insert(tk.END, line_content + "\n")
        result_text.update_idletasks()

        line_start = result_text.index("insert -1 lines")
        line_end = result_text.index("insert")

        pos = line_start
        while True:
            pos = result_text.search(keyword, pos, stopindex=line_end, nocase=True)
            if not pos:
                break
            match_end = f"{pos}+{len(keyword)}c"
            result_text.tag_add("highlight", pos, match_end)
            pos = match_end

    result_text.insert(tk.END, "═" * 80 + "\n", "separator")
    result_text.see(tk.END)

# UI event
def browse_folder():
    folder_path.set(filedialog.askdirectory())

def toggle_search():
    global search_cancelled, is_searching
    if not is_searching:
        # Start search
        search_cancelled = False
        is_searching = True
        search_button.config(text="Cancel")
        status_text.set("Searching...")
        threading.Thread(target=perform_search, daemon=True).start()
    else:
        # Cancel search
        search_cancelled = True
        status_text.set("Cancelling...")
        search_button.config(state="disabled")

def perform_search():
    global is_searching
    folder = folder_path.get().strip()
    keyword = keyword_entry.get().strip()

    if not folder or not keyword:
        messagebox.showwarning("⚠️ Missing Input", "Please select a folder and enter a keyword.")
        search_button.config(text="Search", state="normal")
        status_text.set("Ready")
        is_searching = False
        return

    result_text.delete(1.0, tk.END)
    results = search_keyword_in_folder(folder, keyword, on_result=handle_result)

    result_text.insert(tk.END, f"✅ Search complete. Found {len(results)} matching files.\n")
    search_button.config(text="Search", state="normal")
    is_searching = False
    status_text.set("Search cancelled." if search_cancelled else "Done.")

# GUI
root = tk.Tk()
root.title("📄 Document Search Tool (PDF, Word, PPT, TXT)")
root.geometry("850x600")
root.resizable(True, True)

folder_path = tk.StringVar()
status_text = tk.StringVar(value="Ready")
search_cancelled = False
is_searching = False

main_frame = ttk.Frame(root, padding=15)
main_frame.pack(fill="both", expand=True)

for col in range(8):
    main_frame.columnconfigure(col, weight=0)

main_frame.columnconfigure(1, weight=3)
main_frame.columnconfigure(4, weight=2)
main_frame.columnconfigure(6, weight=1)

main_frame.rowconfigure(3, weight=1)

# Top input row
ttk.Label(main_frame, text="📂 Folder:").grid(row=0, column=0, sticky="e", padx=(0, 5))
folder_entry = ttk.Entry(main_frame, textvariable=folder_path)
folder_entry.grid(row=0, column=1, sticky="we", padx=2)
ttk.Button(main_frame, text="Browse", command=browse_folder).grid(row=0, column=2, padx=(5, 15))

ttk.Label(main_frame, text="🔍 Keyword:").grid(row=0, column=3, sticky="e", padx=(0, 5))
keyword_entry = ttk.Entry(main_frame)
keyword_entry.grid(row=0, column=4, sticky="we", padx=2)
search_button = ttk.Button(main_frame, text="Search", command=toggle_search)
search_button.grid(row=0, column=5, padx=(5, 10))

# Separator
ttk.Separator(main_frame, orient="horizontal").grid(row=2, column=0, columnspan=7, pady=10, sticky="ew")

# Result area with scrollbar
text_frame = tk.Frame(main_frame)
text_frame.grid(row=3, column=0, columnspan=7, sticky="nsew")

result_text = tk.Text(text_frame, wrap="word", borderwidth=1, relief="solid")
result_text.pack(side="left", fill="both", expand=True)
result_text.configure(font=("Courier New", 12))

# Styles
result_text.tag_config("highlight", background="#FFEB3B", foreground="black")
result_text.tag_config("separator", foreground="#708090")
result_text.tag_config("greenbar", background="#B9FBC0", foreground="black")

scrollbar = ttk.Scrollbar(text_frame, orient="vertical", command=result_text.yview)
scrollbar.pack(side="right", fill="y")
result_text.config(yscrollcommand=scrollbar.set)

# Status bar
status_bar = ttk.Label(root, textvariable=status_text, anchor="w", padding=5)
status_bar.pack(side="bottom", fill="x")

root.mainloop()